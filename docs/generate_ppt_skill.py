#!/usr/bin/env python3
"""Generate McKinsey-style PPT using mck-ppt-design skill specifications."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import zipfile, os
from lxml import etree

# ── Color Constants (mck-ppt-design skill) ──
NAVY       = RGBColor(0x05, 0x1C, 0x2C)
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MED_GRAY   = RGBColor(0x66, 0x66, 0x66)
LINE_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
BG_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)

ACCENT_BLUE   = RGBColor(0x00, 0x6B, 0xA6)
ACCENT_GREEN  = RGBColor(0x00, 0x7A, 0x53)
ACCENT_ORANGE = RGBColor(0xD4, 0x6A, 0x00)
ACCENT_RED    = RGBColor(0xC6, 0x28, 0x28)
LIGHT_BLUE    = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_GREEN   = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_ORANGE  = RGBColor(0xFF, 0xF3, 0xE0)
LIGHT_RED     = RGBColor(0xFF, 0xEB, 0xEE)

# ── Typography Constants ──
TITLE_SIZE = Pt(22)
SUB_HEADER_SIZE = Pt(18)
BODY_SIZE = Pt(14)
LABEL_SIZE = Pt(14)
SMALL_SIZE = Pt(9)

# ── Layout Constants ──
SW = Inches(13.333)
SH = Inches(7.5)
LM = Inches(0.8)
CW = Inches(11.733)
BL = None  # will set after prs creation
TOTAL = 14

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BL = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════
# Helper Functions (from mck-ppt-design skill)
# ══════════════════════════════════════════════════════════

def _clean_shape(shape):
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)


def set_ea_font(run, typeface='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)


def add_text(slide, left, top, width, height, text, font_size=Pt(14),
             font_name='Arial', font_color=DARK_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, ea_font='KaiTi', anchor=MSO_ANCHOR.TOP,
             line_spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns','tIns','rIns','bIns']:
        bodyPr.set(attr, '45720')
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
        p.space_before = line_spacing if i > 0 else Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size.pt * 1.35)
        for run in p.runs:
            set_ea_font(run, ea_font)
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    _clean_shape(shape)
    return shape


def add_hline(slide, x, y, length, color=BLACK, thickness=Pt(0.5)):
    h = max(int(thickness), Emu(6350))
    return add_rect(slide, x, y, length, h, color)


def add_oval(slide, x, y, letter, size=Inches(0.45),
             bg=NAVY, fg=WHITE):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = bg
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = letter
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.color.rgb = fg
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        set_ea_font(run, 'KaiTi')
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    _clean_shape(c)
    return c


def add_action_title(slide, text, title_size=Pt(22)):
    add_text(slide, Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.9),
             text, font_size=title_size, font_color=BLACK, bold=True,
             font_name='Georgia', ea_font='KaiTi', anchor=MSO_ANCHOR.MIDDLE)
    add_hline(slide, Inches(0.8), Inches(1.05), Inches(11.7), BLACK, Pt(0.5))


def add_source(slide, text, y=Inches(7.05)):
    add_text(slide, Inches(0.8), y, Inches(11), Inches(0.3),
             text, font_size=SMALL_SIZE, font_color=MED_GRAY)


def add_page_number(slide, num, total=TOTAL):
    add_text(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
             f"{num}/{total}", font_size=SMALL_SIZE, font_color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)


def full_cleanup(outpath):
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Cover (Pattern #1)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_rect(s, 0, 0, SW, Inches(0.05), NAVY)
add_text(s, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
         '微软 Fabric IQ 本体管理', font_size=Pt(44), font_name='Georgia',
         font_color=NAVY, bold=True, ea_font='KaiTi')
add_text(s, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
         '架构详解、构建流程与开源替代方案', font_size=Pt(24),
         font_color=DARK_GRAY, ea_font='KaiTi')
add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
         '技术分享  |  2026年3月', font_size=BODY_SIZE,
         font_color=MED_GRAY, ea_font='KaiTi')
add_hline(s, Inches(1), Inches(6.8), Inches(4), NAVY, Pt(2))

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Table of Contents (Pattern #6)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '目录')
items = [
    ('1', '为什么需要本体', '数据孤岛的根源与语义统一的价值'),
    ('2', 'Fabric IQ 本体架构', '核心概念、构建流程与数据绑定机制'),
    ('3', '开源替代方案', '对标 Fabric IQ 的开源技术栈'),
    ('4', '关键要点与行动建议', '从概念到落地的核心总结'),
]
iy = Inches(1.6)
for num, title, desc in items:
    add_oval(s, LM, iy, num)
    add_text(s, LM + Inches(0.7), iy, Inches(4.0), Inches(0.4),
             title, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
    add_text(s, Inches(5.5), iy + Inches(0.05), Inches(6.5), Inches(0.4),
             desc, font_size=BODY_SIZE, font_color=MED_GRAY)
    iy += Inches(0.7)
    add_hline(s, LM, iy, CW, LINE_GRAY)
    iy += Inches(0.3)
add_source(s, 'Source: 微软本体管理功能深度解析, Fabric IQ 博客')
add_page_number(s, 2)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Executive Summary (Pattern #24 / Three-Stat)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '本体是企业数据与 AI 之间的缺失环节')

# Three metric cards
stats = [
    ('问题', '数据可查询\n但不可理解', ACCENT_RED, LIGHT_RED),
    ('方案', 'Fabric IQ\n语义枢纽', ACCENT_GREEN, LIGHT_GREEN),
    ('替代', 'Protégé + Morph-KGC\n+ Neo4j', ACCENT_BLUE, LIGHT_BLUE),
]
sw = Inches(3.5)
sg = (CW - sw * 3) / 2
for i, (label, desc, accent, bg) in enumerate(stats):
    sx = LM + (sw + sg) * i
    add_rect(s, sx, Inches(1.5), sw, Inches(0.5), accent)
    add_text(s, sx + Inches(0.15), Inches(1.5), sw - Inches(0.3), Inches(0.5),
             label, font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
    add_rect(s, sx, Inches(2.0), sw, Inches(1.5), bg)
    add_text(s, sx + Inches(0.15), Inches(2.1), sw - Inches(0.3), Inches(1.3),
             desc, font_size=BODY_SIZE, font_color=DARK_GRAY,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom insight
add_rect(s, LM, Inches(4.2), CW, Inches(2.2), BG_GRAY)
add_text(s, LM + Inches(0.3), Inches(4.3), Inches(2), Inches(0.35),
         '关键洞见', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
insights = [
    '• 同词异义、数据孤岛阻碍 AI 落地——企业数据缺少统一语义层',
    '• Fabric IQ 通过本体下连数据（OneLake），上接 AI（Data Agent）',
    '• 开源方案 Protégé + Morph-KGC + Neo4j 可复现核心能力',
    '• 本体让数据从"可查"升级为"可理解、可推理、可行动"',
]
add_text(s, LM + Inches(0.3), Inches(4.75), CW - Inches(0.6), Inches(1.5),
         insights, font_size=BODY_SIZE, font_color=DARK_GRAY, line_spacing=Pt(8))

add_source(s, 'Source: 微软本体管理功能深度解析; 开源本体数据映射工具梳理')
add_page_number(s, 3)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture: Three-Layer (Pattern #14 Three-Pillar)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, 'Fabric IQ 本体 = 连接数据与 AI 的语义枢纽')

# Three-layer architecture using horizontal bands
layers = [
    ('应用层', ['Data Agent — 自然语言交互', '知识图谱 — 图遍历查询', '统一语义层 — 共享复用'],
     ACCENT_GREEN, LIGHT_GREEN, Inches(1.25)),
    ('语义层（本体）', ['实体类型 Entity Type', '关系类型 Relationship', '规则 Rule (OWL)', '数据绑定 Data Binding'],
     NAVY, LIGHT_BLUE, Inches(3.25)),
    ('数据层（OneLake）', ['Lakehouse — 批量数据', 'Eventhouse — 实时时序', '语义模型 — Power BI'],
     ACCENT_ORANGE, LIGHT_ORANGE, Inches(5.25)),
]

for layer_name, items, accent, bg, y in layers:
    # Layer header bar
    add_rect(s, LM, y, CW, Inches(0.4), accent)
    add_text(s, LM + Inches(0.15), y, CW - Inches(0.3), Inches(0.4),
             layer_name, font_size=LABEL_SIZE, font_color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    # Layer content area
    add_rect(s, LM, y + Inches(0.4), CW, Inches(1.3), bg)
    n = len(items)
    iw = (CW - Inches(0.3) * (n + 1)) / n
    for j, item in enumerate(items):
        ix = LM + Inches(0.3) + (iw + Inches(0.3)) * j
        add_rect(s, ix, y + Inches(0.55), iw, Inches(0.9), WHITE)
        add_text(s, ix + Inches(0.1), y + Inches(0.6), iw - Inches(0.2), Inches(0.8),
                 item, font_size=LABEL_SIZE, font_color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Arrows between layers
for arrow_y in [Inches(2.95), Inches(4.95)]:
    add_text(s, Inches(5.5), arrow_y, Inches(2), Inches(0.3),
             '\u25BC  \u25BC  \u25BC', font_size=Pt(16), font_color=MED_GRAY,
             alignment=PP_ALIGN.CENTER)

add_source(s, 'Source: Fabric IQ 本体构建流程详解; 微软本体管理功能深度解析')
add_page_number(s, 4)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Core Concepts Table (Pattern #9 Data Table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '8 个核心概念：类比关系型数据库即可快速理解')

headers = ['本体概念', '英文', '定义', '数据库类比']
col_ws = [Inches(1.8), Inches(2.2), Inches(4.5), Inches(3.2)]
rows_data = [
    ['实体类型', 'Entity Type', '业务概念，可绑定数据表', '表结构（Schema）'],
    ['实体类型键', 'Entity Type Key', '唯一标识', '主键（PK）'],
    ['实体实例', 'Entity Instance', '实体的具体数据', '一行记录（Row）'],
    ['属性', 'Property', '实体特征，可绑定字段', '列（Column）'],
    ['关系类型', 'Relationship Type', '实体间关联定义', '外键约束（FK）'],
    ['关系实例', 'Relationship Instance', '关联的具体数据', 'FK 匹配的记录对'],
    ['数据绑定', 'Data Binding', '本体与数据源的映射', 'SQL 映射规则'],
    ['规则', 'Rule', '基于 OWL 2 RL 的逻辑', '触发器 / 存储过程'],
]

hdr_y = Inches(1.5)
cx = LM
for hdr, cw in zip(headers, col_ws):
    add_text(s, cx, hdr_y, cw, Inches(0.4), hdr,
             font_size=BODY_SIZE, font_color=MED_GRAY, bold=True)
    cx += cw
add_hline(s, LM, Inches(1.95), CW, BLACK, Pt(1.0))

row_h = Inches(0.6)
for ri, row in enumerate(rows_data):
    ry = Inches(2.05) + row_h * ri
    cx = LM
    for val, cw in zip(row, col_ws):
        add_text(s, cx, ry, cw, row_h, val, font_size=BODY_SIZE, font_color=DARK_GRAY)
        cx += cw
    add_hline(s, LM, ry + row_h, CW, LINE_GRAY)

# Bottom note
add_rect(s, LM, Inches(6.5), CW, Inches(0.4), BG_GRAY)
add_text(s, LM + Inches(0.2), Inches(6.5), CW - Inches(0.4), Inches(0.4),
         '关键理解：本体中的每个概念都可以直接映射到你熟悉的关系型数据库概念',
         font_size=LABEL_SIZE, font_color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

add_source(s, 'Source: 微软本体管理功能深度解析 — 核心概念表')
add_page_number(s, 5)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Two Paths Comparison (Pattern #19 Side-by-Side)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '两条创建路径：已有模型可秒级生成，新场景需手动构建')

# Left column: Auto
lx = LM
lw = Inches(5.5)
add_rect(s, lx, Inches(1.4), lw, Inches(0.45), ACCENT_GREEN)
add_text(s, lx + Inches(0.15), Inches(1.4), lw - Inches(0.3), Inches(0.45),
         '从语义模型自动生成', font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
auto_points = [
    '• 前提：已有 Power BI 语义模型',
    '• 表 → 实体类型，字段 → 属性（自动）',
    '• 耗时：分钟级',
    '• 灵活性：受限于已有模型结构',
    '• 适用：快速验证、已有资产复用',
]
add_text(s, lx + Inches(0.15), Inches(2.05), lw - Inches(0.3), Inches(3.0),
         auto_points, font_size=BODY_SIZE, font_color=DARK_GRAY, line_spacing=Pt(10))

# Right column: Manual
rx = LM + Inches(6.2)
rw = Inches(5.5)
add_rect(s, rx, Inches(1.4), rw, Inches(0.45), NAVY)
add_text(s, rx + Inches(0.15), Inches(1.4), rw - Inches(0.3), Inches(0.45),
         '手动构建', font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
manual_points = [
    '• 前提：无',
    '• 逐步定义实体、属性、关系',
    '• 耗时：天级',
    '• 灵活性：完全自定义',
    '• 适用：全新业务建模、深度定制',
]
add_text(s, rx + Inches(0.15), Inches(2.05), rw - Inches(0.3), Inches(3.0),
         manual_points, font_size=BODY_SIZE, font_color=DARK_GRAY, line_spacing=Pt(10))

# Divider line between columns
add_hline(s, LM + Inches(5.8), Inches(1.4), Pt(0.5), LINE_GRAY)

# Bottom recommendation
add_rect(s, LM, Inches(5.3), CW, Inches(0.9), BG_GRAY)
add_text(s, LM + Inches(0.3), Inches(5.35), Inches(2), Inches(0.3),
         '实践建议', font_size=LABEL_SIZE, font_color=NAVY, bold=True)
add_text(s, LM + Inches(0.3), Inches(5.7), CW - Inches(0.6), Inches(0.4),
         '先用手动构建理解核心概念，再用自动生成快速扩展。两种方式可混合使用。',
         font_size=BODY_SIZE, font_color=DARK_GRAY)

add_source(s, 'Source: Fabric IQ 本体构建流程详解 — 两种创建方式')
add_page_number(s, 6)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — 5-Phase Build Flow (Pattern #16 Process Chevron)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '5 阶段构建流程：从空白本体到可查询的知识图谱')

phases = [
    ('01', '创建本体', '新建 Ontology 项目\n选择创建方式'),
    ('02', '定义实体类型', '添加 Entity Type\n设置属性与主键\n绑定数据源'),
    ('03', '定义关系', '连接实体类型\n描述业务关联\n绑定关系数据'),
    ('04', '配置规则', 'OWL 2 RL 规则逻辑\n触发自动化动作'),
    ('05', '预览验证', '查看实体实例\n可视化关系图\n图引擎查询测试'),
]
pw = Inches(2.1)
pg = (CW - pw * 5) / 4
colors = [NAVY, ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]

for i, (num, title, desc) in enumerate(phases):
    px = LM + (pw + pg) * i
    # Number circle
    add_oval(s, px, Inches(1.4), num, size=Inches(0.4), bg=colors[i])
    # Title
    add_text(s, px + Inches(0.5), Inches(1.4), Inches(1.6), Inches(0.4),
             title, font_size=LABEL_SIZE, font_color=colors[i], bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    # Description box
    add_rect(s, px, Inches(1.95), pw, Inches(2.0), BG_GRAY)
    add_text(s, px + Inches(0.1), Inches(2.05), pw - Inches(0.2), Inches(1.8),
             desc, font_size=LABEL_SIZE, font_color=DARK_GRAY, line_spacing=Pt(6))
    # Arrow
    if i < 4:
        add_text(s, px + pw, Inches(2.2), pg, Inches(0.4),
                 '\u25B6', font_size=Pt(18), font_color=colors[i],
                 alignment=PP_ALIGN.CENTER)

# Output section
add_hline(s, LM, Inches(4.25), CW, LINE_GRAY)
add_text(s, LM, Inches(4.4), CW, Inches(0.3),
         '构建产出', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)

outputs = [
    ('知识图谱 (Graph)', '自动生成，支持图遍历查询', ACCENT_GREEN),
    ('Data Agent', 'AI 智能体，自然语言交互', ACCENT_BLUE),
    ('统一语义层', '下游应用共享业务语义', NAVY),
]
ow = Inches(3.5)
og = (CW - ow * 3) / 2
for i, (otitle, odesc, ocolor) in enumerate(outputs):
    ox = LM + (ow + og) * i
    add_rect(s, ox, Inches(4.85), ow, Inches(1.6), WHITE)
    add_hline(s, ox, Inches(4.85), ow, ocolor, Pt(3))
    add_text(s, ox + Inches(0.15), Inches(5.0), ow - Inches(0.3), Inches(0.4),
             otitle, font_size=LABEL_SIZE, font_color=ocolor, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s, ox + Inches(0.15), Inches(5.45), ow - Inches(0.3), Inches(0.8),
             odesc, font_size=LABEL_SIZE, font_color=MED_GRAY,
             alignment=PP_ALIGN.CENTER)

add_source(s, 'Source: Fabric IQ 本体构建流程详解 — 完整构建流程')
add_page_number(s, 7)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Data Binding Deep Dive (Pattern #11 Data Table + Insight)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '数据绑定机制：本体不存数据，只做"翻译"')

# Example table
ex_headers = ['实体类型', '属性', '主键', '绑定表', '数据源']
ex_col_ws = [Inches(2.2), Inches(3.3), Inches(1.3), Inches(2.2), Inches(2.7)]
ex_rows = [
    ['Store（门店）', 'store_id, name, address, region', 'store_id', 'dim_store', 'Lakehouse'],
    ['Employee（员工）', 'emp_id, name, role, hire_date', 'emp_id', 'dim_employee', 'Lakehouse'],
    ['SaleEvent（销售事件）', 'event_id, timestamp, amount', 'event_id', 'sales_stream', 'Eventhouse'],
]

hdr_y = Inches(1.4)
cx = LM
for hdr, cw in zip(ex_headers, ex_col_ws):
    add_text(s, cx, hdr_y, cw, Inches(0.35), hdr,
             font_size=LABEL_SIZE, font_color=MED_GRAY, bold=True)
    cx += cw
add_hline(s, LM, Inches(1.8), CW, BLACK, Pt(1.0))

for ri, row in enumerate(ex_rows):
    ry = Inches(1.9) + Inches(0.5) * ri
    cx = LM
    for val, cw in zip(row, ex_col_ws):
        add_text(s, cx, ry, cw, Inches(0.5), val, font_size=LABEL_SIZE, font_color=DARK_GRAY)
        cx += cw
    add_hline(s, LM, ry + Inches(0.5), CW, LINE_GRAY)

# Mapping rules table (left)
add_text(s, LM, Inches(3.7), Inches(5), Inches(0.3),
         '映射规则', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)

map_data = [
    ['本体层（语义）', '\u2192', '数据层（存储）'],
    ['Entity Type', '\u2192', 'Table（表）'],
    ['Property', '\u2192', 'Column（字段）'],
    ['Entity Type Key', '\u2192', 'Primary Key（主键）'],
    ['Relationship', '\u2192', 'Join Table / FK'],
]
map_cws = [Inches(2.0), Inches(0.5), Inches(2.5)]
for ri, row in enumerate(map_data):
    ry = Inches(4.1) + Inches(0.4) * ri
    cx = LM
    for val, cw in zip(row, map_cws):
        fc = NAVY if ri == 0 else DARK_GRAY
        fb = ri == 0
        add_text(s, cx, ry, cw, Inches(0.35), val, font_size=LABEL_SIZE,
                 font_color=fc, bold=fb)
        cx += cw
    if ri == 0:
        add_hline(s, LM, Inches(4.45), Inches(5), BLACK, Pt(1.0))
    elif ri < len(map_data):
        add_hline(s, LM, ry + Inches(0.4), Inches(5), LINE_GRAY)

# How it works (right)
add_rect(s, Inches(6.8), Inches(3.7), Inches(5.7), Inches(2.8), BG_GRAY)
add_text(s, Inches(7.1), Inches(3.8), Inches(5.2), Inches(0.35),
         '工作原理', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
steps = [
    '1. 用户在 Protégé / Fabric 中定义本体结构',
    '2. 通过 Data Binding 映射到底层表和字段',
    '3. 查询时，系统自动翻译为 SQL / KQL',
    '4. 结果组装为 Entity Instance 返回',
    '',
    '本体 = 元数据映射层，不复制数据',
]
add_text(s, Inches(7.1), Inches(4.25), Inches(5.2), Inches(2.0),
         steps, font_size=LABEL_SIZE, font_color=DARK_GRAY, line_spacing=Pt(6))

add_source(s, 'Source: Fabric IQ 本体构建流程详解 — 数据绑定')
add_page_number(s, 8)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Relationships & Rules (Pattern #33 Case Study)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '关系与规则：让数据从"可查"升级为"可推理"')

# Left: Relationship diagram
add_text(s, LM, Inches(1.3), Inches(5.5), Inches(0.35),
         '关系类型（Relationship Type）', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)

# Entity nodes
entities = [
    ('Store\n门店', LM, Inches(2.3)),
    ('Employee\n员工', LM + Inches(3.2), Inches(2.3)),
    ('SaleEvent\n销售事件', LM + Inches(3.2), Inches(4.0)),
]
for label, ex, ey in entities:
    add_rect(s, ex, ey, Inches(2.0), Inches(0.8), NAVY)
    add_text(s, ex + Inches(0.1), ey, Inches(1.8), Inches(0.8),
             label, font_size=LABEL_SIZE, font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Relationship labels
add_text(s, LM + Inches(2.0), Inches(2.0), Inches(1.5), Inches(0.3),
         'employs \u25B6', font_size=SMALL_SIZE, font_color=MED_GRAY,
         alignment=PP_ALIGN.CENTER)
add_text(s, LM + Inches(1.2), Inches(3.4), Inches(2.5), Inches(0.3),
         'generates \u25B6', font_size=SMALL_SIZE, font_color=MED_GRAY,
         alignment=PP_ALIGN.CENTER)
add_text(s, LM + Inches(5.0), Inches(3.1), Inches(1.5), Inches(0.3),
         'completes \u25B2', font_size=SMALL_SIZE, font_color=MED_GRAY,
         alignment=PP_ALIGN.CENTER)

# Right: Rules
add_text(s, Inches(7.2), Inches(1.3), Inches(5), Inches(0.35),
         '规则引擎（OWL 2 RL）', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)

add_rect(s, Inches(7.2), Inches(1.8), Inches(5.3), Inches(2.0), BG_GRAY)
rule_lines = [
    '规则示例',
    'IF:',
    '  SaleEvent.amount < threshold',
    '  AND Store.region = "华东"',
    'THEN:',
    '  触发告警动作 (Action)',
]
add_text(s, Inches(7.4), Inches(1.85), Inches(5.0), Inches(1.8),
         rule_lines, font_size=LABEL_SIZE, font_color=DARK_GRAY,
         font_name='Consolas', line_spacing=Pt(4))

# Bottom: Key capabilities
add_hline(s, LM, Inches(5.2), CW, LINE_GRAY)
add_text(s, LM, Inches(5.35), CW, Inches(0.3),
         '核心能力', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)

caps = [
    ('关系 = 实体间的业务关联', '支持多对多，可绑定中间表\n关系实例随数据自动生成'),
    ('规则 = 业务逻辑的形式化', '基于 W3C OWL 2 RL 标准\n能触发告警、通知、API 调用'),
    ('推理 = 从已知推导未知', '基于实体属性和关系\n自动发现隐含的业务洞察'),
]
cw = (CW - Inches(0.2) * 2) / 3
for i, (ctitle, cdesc) in enumerate(caps):
    cx = LM + (cw + Inches(0.2)) * i
    add_rect(s, cx, Inches(5.75), cw, Inches(1.2), BG_GRAY)
    add_hline(s, cx, Inches(5.75), cw, ACCENT_BLUE, Pt(2))
    add_text(s, cx + Inches(0.1), Inches(5.85), cw - Inches(0.2), Inches(0.3),
             ctitle, font_size=LABEL_SIZE, font_color=NAVY, bold=True)
    add_text(s, cx + Inches(0.1), Inches(6.2), cw - Inches(0.2), Inches(0.6),
             cdesc, font_size=SMALL_SIZE, font_color=MED_GRAY)

add_source(s, 'Source: 微软本体管理功能深度解析 — 规则引擎; Fabric IQ 本体构建流程详解')
add_page_number(s, 9)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Three Access Methods (Pattern #12 Metric Cards)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '三种访问方式：从技术到业务的完整覆盖')

methods = [
    ('A', '本体实例查询', ['原理：直接查询实体和关系实例', '语言：类 SQL',
     '人群：数据工程师', '场景：数据校验、实例浏览', '门槛：\u2605\u2605']),
    ('B', '图引擎查询', ['原理：图遍历，路径发现', '语言：Cypher / SPARQL',
     '人群：图分析师', '场景：关联分析、社区发现', '门槛：\u2605\u2605\u2605']),
    ('C', '自然语言对话', ['原理：Data Agent 翻译自然语言', '语言：自然语言',
     '人群：业务人员 / 管理层', '场景：业务问答、数据探索', '门槛：\u2605']),
]
cw = Inches(3.5)
cg = (CW - cw * 3) / 2
colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE]

for i, (letter, title, points) in enumerate(methods):
    cx = LM + (cw + cg) * i
    # Card background
    add_rect(s, cx, Inches(1.4), cw, Inches(4.2), BG_GRAY)
    # Top accent bar
    add_rect(s, cx, Inches(1.4), cw, Inches(0.06), colors[i])
    # Circle label
    add_oval(s, cx + cw // 2 - Inches(0.22), Inches(1.6), letter,
             size=Inches(0.45), bg=colors[i])
    # Title
    add_text(s, cx + Inches(0.15), Inches(2.15), cw - Inches(0.3), Inches(0.35),
             title, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_hline(s, cx + Inches(0.3), Inches(2.55), cw - Inches(0.6), LINE_GRAY)
    # Points
    add_text(s, cx + Inches(0.2), Inches(2.7), cw - Inches(0.4), Inches(2.5),
             [f'• {p}' for p in points], font_size=LABEL_SIZE, font_color=DARK_GRAY,
             line_spacing=Pt(8))

# Bottom insight
add_rect(s, LM, Inches(5.85), CW, Inches(0.7), BG_GRAY)
add_text(s, LM + Inches(0.3), Inches(5.9), CW - Inches(0.6), Inches(0.6),
         '三种方式共享同一套本体定义，确保"技术查询"和"业务提问"得到语义一致的结果。本体是唯一的语义真相来源。',
         font_size=LABEL_SIZE, font_color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

add_source(s, 'Source: 微软本体管理功能深度解析 — 三种数据访问方式')
add_page_number(s, 10)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Open Source Alternatives (Pattern #11 Data Table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '开源替代方案：每个组件都有成熟的开源选择')

oss_headers = ['Fabric IQ 功能', '开源方案', '工具 / 标准', '说明']
oss_col_ws = [Inches(2.2), Inches(2.5), Inches(3.2), Inches(3.8)]
oss_rows = [
    ['本体编辑', 'Protégé', 'OWL / RDF', '最成熟的开源本体编辑器'],
    ['数据绑定', 'R2RML / RML', 'W3C 标准映射语言', '表→实体，字段→属性'],
    ['映射执行', 'Morph-KGC', 'Python, R2RML+RML', '首选引擎，支持大数据量'],
    ['图存储', 'Neo4j 社区版', 'Cypher 查询', '图数据库事实标准'],
    ['图存储', 'Apache Jena', 'SPARQL + RDF', '语义 Web 标准栈'],
    ['规则推理', 'GraphDB / Jena', 'OWL 推理引擎', '内置 OWL 2 RL 推理'],
    ['本体桥接', 'Neosemantics', 'Neo4j 插件', 'RDF/OWL 导入导出'],
    ['自然语言', 'LangChain + Neo4j', 'GraphRAG', 'LLM + 图查询'],
]

hdr_y = Inches(1.4)
cx = LM
for hdr, cw in zip(oss_headers, oss_col_ws):
    add_text(s, cx, hdr_y, cw, Inches(0.35), hdr,
             font_size=LABEL_SIZE, font_color=MED_GRAY, bold=True)
    cx += cw
add_hline(s, LM, Inches(1.8), CW, BLACK, Pt(1.0))

row_h = Inches(0.55)
for ri, row in enumerate(oss_rows):
    ry = Inches(1.9) + row_h * ri
    cx = LM
    for val, cw in zip(row, oss_col_ws):
        add_text(s, cx, ry, cw, row_h, val, font_size=LABEL_SIZE, font_color=DARK_GRAY)
        cx += cw
    add_hline(s, LM, ry + row_h, CW, LINE_GRAY)

# Bottom note
add_rect(s, LM, Inches(6.4), CW, Inches(0.4), BG_GRAY)
add_text(s, LM + Inches(0.2), Inches(6.4), CW - Inches(0.4), Inches(0.4),
         '基于 W3C 开放标准（OWL、RDF、R2RML、RML），所有组件可自由组合，无供应商锁定',
         font_size=LABEL_SIZE, font_color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

add_source(s, 'Source: 开源本体数据映射工具梳理')
add_page_number(s, 11)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — Recommended Stack (Pattern #16 Process Chevron)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '推荐最小可行技术栈：3 个工具覆盖完整链路')

# Pipeline steps
steps = [
    ('1', 'Protégé', '定义本体\n(OWL/RDF)', NAVY),
    ('2', 'YARRRML', '编写映射规则\n(表→实体)', ACCENT_BLUE),
    ('3', 'Morph-KGC', '执行映射\n生成知识图谱', ACCENT_GREEN),
    ('4', 'Neo4j', '存储查询\n图遍历分析', ACCENT_ORANGE),
]
sw = Inches(2.5)
sg = (CW - sw * 4) / 3
for i, (num, title, desc, color) in enumerate(steps):
    sx = LM + (sw + sg) * i
    add_rect(s, sx, Inches(1.5), sw, Inches(0.5), color)
    add_text(s, sx + Inches(0.15), Inches(1.5), sw - Inches(0.3), Inches(0.5),
             f'Step {num}', font_size=LABEL_SIZE, font_color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
    add_rect(s, sx, Inches(2.0), sw, Inches(1.2), BG_GRAY)
    add_text(s, sx + Inches(0.15), Inches(2.1), sw - Inches(0.3), Inches(0.35),
             title, font_size=SUB_HEADER_SIZE, font_color=color, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s, sx + Inches(0.15), Inches(2.5), sw - Inches(0.3), Inches(0.6),
             desc, font_size=LABEL_SIZE, font_color=DARK_GRAY,
             alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, sx + sw, Inches(1.6), sg, Inches(0.4),
                 '\u25B6', font_size=Pt(18), font_color=color,
                 alignment=PP_ALIGN.CENTER)

# Why this stack
add_hline(s, LM, Inches(3.55), CW, LINE_GRAY)
add_text(s, LM, Inches(3.7), CW, Inches(0.3),
         '选型理由', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)

reasons = [
    ('Protégé', '全球使用最广泛的本体编辑器，学术+企业社区成熟，支持 OWL 2 全特性'),
    ('Morph-KGC', 'Python 原生，支持 R2RML/RML/RML-star，可扩展到千万级三元组'),
    ('Neo4j', '图数据库事实标准，Cypher 查询直观，APOC/GraphDataScience 生态丰富'),
    ('整体', '基于 W3C 开放标准，组件可替换，无供应商锁定风险'),
]
for i, (tool, reason) in enumerate(reasons):
    ry = Inches(4.15) + Inches(0.55) * i
    add_oval(s, LM, ry + Inches(0.05), str(i + 1), size=Inches(0.35))
    add_text(s, LM + Inches(0.55), ry, Inches(1.5), Inches(0.45),
             tool, font_size=LABEL_SIZE, font_color=NAVY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, LM + Inches(2.0), ry, Inches(9.5), Inches(0.45),
             reason, font_size=LABEL_SIZE, font_color=DARK_GRAY,
             anchor=MSO_ANCHOR.MIDDLE)

add_source(s, 'Source: 开源本体数据映射工具梳理 — 推荐组合')
add_page_number(s, 12)

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — Key Takeaways (Pattern #25 Key Takeaway)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '关键要点')

takeaways = [
    ('1', '本体是语义枢纽，不是数据存储',
     '通过 Data Binding 映射底层数据，自身只维护元数据定义——本体 = 业务的"词典" + "语法"'),
    ('2', '构建有规律：实体 → 数据绑定 → 关系 → 规则',
     '5 阶段流程可复用，建议先用自动生成跑通，再手动精调'),
    ('3', '开源方案已能覆盖完整链路',
     'Protégé + Morph-KGC + Neo4j 是最小可行组合，基于 W3C 标准无锁定'),
    ('4', '最终目标：让数据可理解、可推理、可行动',
     '从 SQL 查询到自然语言对话，本体是连接企业数据与 AI 的桥梁'),
]

for i, (num, title, desc) in enumerate(takeaways):
    ty = Inches(1.4) + Inches(1.35) * i
    add_oval(s, LM, ty, num, size=Inches(0.45))
    add_text(s, LM + Inches(0.65), ty, Inches(10.5), Inches(0.35),
             title, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, LM + Inches(0.65), ty + Inches(0.45), Inches(10.5), Inches(0.5),
             desc, font_size=LABEL_SIZE, font_color=MED_GRAY)

# Bottom quote
add_rect(s, LM, Inches(6.2), CW, Inches(0.7), NAVY)
add_text(s, LM + Inches(0.3), Inches(6.25), CW - Inches(0.6), Inches(0.6),
         '"仅靠数据模型无法做好 Context 的落地，唯有引入本体，才能让数据真正变得可理解、可推理、可应用。"',
         font_size=LABEL_SIZE, font_color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_source(s, 'Source: 微软本体管理功能深度解析 — 关键结论')
add_page_number(s, 13)

# ══════════════════════════════════════════════════════════════
# SLIDE 14 — References (Pattern #11 Data Table)
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BL)
add_action_title(s, '附录：参考资源与延伸阅读')

ref_headers = ['分类', '资源', '链接']
ref_col_ws = [Inches(1.2), Inches(4.5), Inches(6.0)]
ref_rows = [
    ['官方', 'Fabric IQ 本体教程 (4步)', 'learn.microsoft.com/fabric/iq/ontology'],
    ['官方', 'Fabric IQ 博客: Semantic Foundation for AI', 'blog.fabric.microsoft.com'],
    ['工具', 'Protégé 本体编辑器', 'protege.stanford.edu'],
    ['工具', 'Morph-KGC 映射引擎', 'github.com/morph-kgc/morph-kgc'],
    ['工具', 'Neo4j 图数据库', 'neo4j.com'],
    ['工具', 'Neosemantics (n10s) Neo4j 插件', 'neo4j.com/blog/nodes/neosemantics'],
    ['标准', 'W3C R2RML 映射语言', 'w3.org/TR/r2rml'],
    ['标准', 'RML 映射语言', 'rml.io/specs/rml'],
    ['标准', 'OWL 2 Web Ontology Language', 'w3.org/TR/owl2-overview'],
    ['社区', 'Awesome KGC Tools 列表', 'kg-construct.github.io/awesome-kgc-tools'],
    ['社区', 'OpenKG 中文开放知识图谱', 'openkg.cn'],
]

hdr_y = Inches(1.4)
cx = LM
for hdr, cw in zip(ref_headers, ref_col_ws):
    add_text(s, cx, hdr_y, cw, Inches(0.35), hdr,
             font_size=LABEL_SIZE, font_color=MED_GRAY, bold=True)
    cx += cw
add_hline(s, LM, Inches(1.8), CW, BLACK, Pt(1.0))

row_h = Inches(0.42)
for ri, row in enumerate(ref_rows):
    ry = Inches(1.9) + row_h * ri
    cx = LM
    for val, cw in zip(row, ref_col_ws):
        add_text(s, cx, ry, cw, row_h, val, font_size=LABEL_SIZE, font_color=DARK_GRAY)
        cx += cw
    add_hline(s, LM, ry + row_h, CW, LINE_GRAY)

add_source(s, 'Source: 综合参考资源')
add_page_number(s, 14)

# ── Save & Cleanup ──
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Fabric-IQ本体管理技术分享-Skill版.pptx")
prs.save(out)
full_cleanup(out)
print(f"Done: {out}")
print(f"Slides: {len(prs.slides)}")
